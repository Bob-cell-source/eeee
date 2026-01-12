# file_parse_tool.py
from __future__ import annotations
import base64, csv, re, tempfile
from typing import Any, Dict, List, Optional, Tuple

import httpx
from .utils import is_url, normalize_url
from ..config import get_settings

MAX_DOWNLOAD_SIZE = get_settings().file_parse_max_download_size

async def _download_to_tmp(url: str, timeout_s: int = 30) -> Tuple[str, str]:
    """Download file with size limit"""
    async with httpx.AsyncClient(timeout=timeout_s, follow_redirects=True) as client:
        async with client.stream("GET", url) as r:
            r.raise_for_status()
            ct = r.headers.get("content-type", "")
            
            # Check content-length if available
            content_length = r.headers.get("content-length")
            if content_length and int(content_length) > MAX_DOWNLOAD_SIZE:
                raise ValueError(f"File too large: {content_length} bytes > {MAX_DOWNLOAD_SIZE} bytes")
            
            suffix = ".bin"
            if "pdf" in ct: suffix = ".pdf"
            elif "html" in ct: suffix = ".html"
            elif "text/plain" in ct: suffix = ".txt"
            
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
            size = 0
            async for chunk in r.aiter_bytes():
                size += len(chunk)
                if size > MAX_DOWNLOAD_SIZE:
                    tmp.close()
                    raise ValueError(f"File too large: exceeds {MAX_DOWNLOAD_SIZE} bytes")
                tmp.write(chunk)
            tmp.close()
            return tmp.name, ct

def _guess_ext(path: str, content_type: Optional[str] = None) -> str:
    p = path.lower()
    for ext in (".pdf",".docx",".pptx",".xlsx",".csv",".md",".txt",".html",".htm",".png",".jpg",".jpeg",".webp"):
        if p.endswith(ext): return ext
    if content_type:
        ct = content_type.lower()
        if "pdf" in ct: return ".pdf"
        if "text/html" in ct: return ".html"
        if "text/plain" in ct: return ".txt"
    return ""

def _read_text(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def _parse_html_basic(path: str) -> str:
    html = _read_text(path)
    html = re.sub(r"(?is)<(script|style|noscript).*?>.*?</\1>", " ", html)
    text = re.sub(r"(?s)<[^>]+>", " ", html)
    return re.sub(r"\s+", " ", text).strip()

def _parse_csv(path: str, max_rows: int = 200) -> str:
    out=[]
    with open(path, "r", encoding="utf-8", errors="ignore", newline="") as f:
        reader=csv.reader(f)
        for i,row in enumerate(reader, start=1):
            if i>max_rows: out.append("... (truncated)"); break
            out.append("\t".join(row))
    return "\n".join(out)

def _parse_xlsx(path: str, max_rows: int = 200) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    out=[]
    for sh in wb.worksheets:
        out.append(f"# Sheet: {sh.title}")
        for i,row in enumerate(sh.iter_rows(values_only=True), start=1):
            if i>max_rows: out.append("... (truncated)"); break
            vals=["" if v is None else str(v) for v in row]
            if any(v.strip() for v in vals): out.append("\t".join(vals))
    return "\n".join(out)

def _parse_docx(path: str) -> str:
    from docx import Document
    doc=Document(path)
    chunks=[]
    for p in doc.paragraphs:
        if p.text.strip(): chunks.append(p.text.strip())
    for table in doc.tables:
        for row in table.rows:
            cells=[c.text.strip() for c in row.cells]
            if any(cells): chunks.append(" | ".join(cells))
    return "\n".join(chunks)

def _parse_pptx(path: str) -> str:
    from pptx import Presentation
    prs=Presentation(path)
    chunks=[]
    for si,slide in enumerate(prs.slides, start=1):
        chunks.append(f"# Slide {si}")
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text and text.strip():
                chunks.append(text.strip())
    return "\n".join(chunks)

def _parse_pdf_text_and_images(path: str, extract_images: bool) -> Tuple[str, List[bytes]]:
    import fitz  # pymupdf
    doc = fitz.open(path)
    texts: List[str] = []
    images: List[bytes] = []

    try:
        for i in range(doc.page_count):
            page: Any = doc[i]  # <- 关键：让 Pylance 不再乱推断
            texts.append(page.get_text("text"))
            if extract_images:
                for img in page.get_images(full=True):
                    xref = img[0]
                    base = doc.extract_image(xref)
                    if base and "image" in base:
                        images.append(base["image"])
    finally:
        doc.close()

    return "\n".join(texts), images

class FileParseClient:
    async def parse(self, file: str, extract_images: bool = False, max_chars: int = 60000, timeout_s: int = 30) -> Dict[str, Any]:
        local_path=file
        ct=None
        if is_url(file):
            local_path, ct = await _download_to_tmp(file, timeout_s=timeout_s)

        ext=_guess_ext(local_path, ct)
        meta={"source": file, "local_path": local_path if is_url(file) else None, "ext": ext}

        images: List[bytes] = []
        try:
            if ext==".pdf":
                text, images = _parse_pdf_text_and_images(local_path, extract_images)
            elif ext==".docx":
                text=_parse_docx(local_path)
            elif ext==".pptx":
                text=_parse_pptx(local_path)
            elif ext==".xlsx":
                text=_parse_xlsx(local_path)
            elif ext==".csv":
                text=_parse_csv(local_path)
            elif ext in (".md",".txt"):
                text=_read_text(local_path)
            elif ext in (".html",".htm"):
                text=_parse_html_basic(local_path)
            else:
                text=_read_text(local_path)
        except ImportError as e:
            missing = str(e).split("'")[1] if "'" in str(e) else "unknown"
            return {"error": f"缺少解析依赖 {missing}，请运行: pip install {missing}", "metadata": meta}
        except ValueError as e:
            # 文件过大等验证错误
            return {"error": str(e), "metadata": meta}
        except Exception as e:
            return {"error": repr(e), "metadata": meta}

        if len(text) > max_chars:
            text = text[:max_chars] + "\n\n... (truncated)"

        # 图片 bytes 转 base64，给集成层或 image_parse 使用
        images_b64 = [base64.b64encode(b).decode("utf-8") for b in images]
        return {"text": text, "metadata": meta, "images_base64": images_b64}


# 单例模式
_file_parser_singleton: Optional[FileParseClient] = None

def get_file_parser() -> FileParseClient:
    """获取文件解析器单例"""
    global _file_parser_singleton
    if _file_parser_singleton is None:
        _file_parser_singleton = FileParseClient()
    return _file_parser_singleton
